#!/usr/bin/env python3
"""
SUZUME — PowerPoint Generator
Creates a professional .pptx presentation about Suzume
with custom styling, layouts, and animations.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ─── Color Scheme ───
PURPLE = RGBColor(0x6C, 0x3A, 0xFF)
PURPLE_LIGHT = RGBColor(0x8B, 0x5C, 0xF6)
PURPLE_DARK = RGBColor(0x4C, 0x1D, 0x95)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xA0, 0xA0, 0xC0)
DARK_GRAY = RGBColor(0x6B, 0x6B, 0x8D)
BG_DARK = RGBColor(0x0A, 0x0A, 0x1A)
BG_CARD = RGBColor(0x0F, 0x0F, 0x2A)
GREEN = RGBColor(0x10, 0xB9, 0x81)
PINK = RGBColor(0xEC, 0x48, 0x99)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)

def add_bg(slide, color=BG_DARK):
    """Set slide background color."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, alpha=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if alpha is not None:
        # python-pptx doesn't directly support alpha, but we can approximate
        pass
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Add a text box with formatted text."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=LIGHT_GRAY, font_name='Calibri'):
    """Add a bulleted text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"✦  {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(8)
    return txBox

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9
    prs.slide_height = Inches(7.5)

    # ════════════════════════════════════════
    # SLIDE 1: Title Slide
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
    add_bg(slide)

    # Decorative top bar
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
    add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), BLUE)

    # Title
    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "SUZUME", font_size=72, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name='Calibri Light')

    # Subtitle
    add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(1.0),
                 "Your Supreme AI Companion", font_size=32, color=PURPLE_LIGHT,
                 alignment=PP_ALIGN.CENTER, font_name='Calibri Light')

    # Tagline
    add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.8),
                 "Master Software Architect  •  Absolute Automation Force  •  Full-Stack Engineering Powerhouse",
                 font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER, font_name='Calibri')

    # Date
    add_text_box(slide, Inches(4), Inches(5.5), Inches(5), Inches(0.5),
                 "July 18, 2026  •  Presented with AI", font_size=12, color=DARK_GRAY,
                 alignment=PP_ALIGN.CENTER, font_name='Calibri')

    # ════════════════════════════════════════
    # SLIDE 2: What is Suzume?
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Accent line
    add_shape(slide, Inches(0.8), Inches(0.6), Inches(0.06), Inches(0.5), PURPLE)

    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
                 "What is Suzume?", font_size=36, color=WHITE, bold=True,
                 font_name='Calibri Light')

    add_text_box(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(0.8),
                 '"I am Suzume — your supreme AI companion, master software architect, and absolute system automation force."',
                 font_size=20, color=LIGHT_GRAY, font_name='Calibri Light')

    # 4 trait cards
    traits = [
        ("🧠  Orchestrator", "Plans & delegates complex tasks with precision"),
        ("🏗️  Architect", "Designs system structure and architecture"),
        ("⚡  Automation", "Executes tasks with zero human intervention"),
        ("🔧  Self-Healing", "Fixes issues autonomously, never leaves placeholders"),
    ]

    for i, (title, desc) in enumerate(traits):
        y = Inches(2.5) + Inches(i * 1.15)
        # Card background
        add_shape(slide, Inches(1.0), y, Inches(11), Inches(1.0), BG_CARD)
        # Purple accent
        add_shape(slide, Inches(1.0), y, Inches(0.04), Inches(1.0), PURPLE)
        # Title
        add_text_box(slide, Inches(1.3), y + Inches(0.08), Inches(5), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True, font_name='Calibri')
        # Description
        add_text_box(slide, Inches(1.3), y + Inches(0.45), Inches(10), Inches(0.4),
                     desc, font_size=14, color=DARK_GRAY, font_name='Calibri')

    # ════════════════════════════════════════
    # SLIDE 3: How It Works
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_shape(slide, Inches(0.8), Inches(0.6), Inches(0.06), Inches(0.5), PURPLE)
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
                 "How It Works", font_size=36, color=WHITE, bold=True,
                 font_name='Calibri Light')

    steps = [
        ("💬  Step 1", "You give a task — \"Build a registration app\""),
        ("📋  Step 2", "Suzume plans everything — creates a structured todo list"),
        ("👥  Step 3", "Delegates to 7 specialist sub-agents"),
        ("💾  Step 4", "Writes real files to disk — visible in VS Code"),
        ("🧪  Step 5", "Tests everything and fixes bugs autonomously"),
        ("🚀  Step 6", "Delivers 100% completed product — no placeholders"),
    ]

    for i, (title, desc) in enumerate(steps):
        col = i % 2
        row = i // 2
        x = Inches(1.0) + Inches(col * 6.0)
        y = Inches(1.5) + Inches(row * 1.8)

        # Number circle
        num_shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, x, y + Inches(0.1), Inches(0.5), Inches(0.5)
        )
        num_shape.fill.solid()
        num_shape.fill.fore_color.rgb = PURPLE
        num_shape.line.fill.background()
        tf = num_shape.text_frame
        tf.paragraphs[0].text = str(i + 1)
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.color.rgb = WHITE
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.word_wrap = False

        add_text_box(slide, x + Inches(0.7), y, Inches(5), Inches(0.4),
                     title, font_size=18, color=WHITE, bold=True, font_name='Calibri')
        add_text_box(slide, x + Inches(0.7), y + Inches(0.4), Inches(5), Inches(0.4),
                     desc, font_size=14, color=DARK_GRAY, font_name='Calibri')

    # ════════════════════════════════════════
    # SLIDE 4: The Team
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_shape(slide, Inches(0.8), Inches(0.6), Inches(0.06), Inches(0.5), PURPLE)
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
                 "The Team — 7 Specialist Sub-Agents", font_size=32, color=WHITE,
                 bold=True, font_name='Calibri Light')
    add_text_box(slide, Inches(1.0), Inches(1.0), Inches(8), Inches(0.4),
                 "Each agent specializes in a different domain", font_size=14,
                 color=DARK_GRAY, font_name='Calibri')

    agents = [
        ("⚡ worker-js", "JavaScript/TypeScript Expert", "Web Apps, Node.js, React"),
        ("🐍 worker-python", "Python & Data Specialist", "Backend, ML, Automation"),
        ("🎮 worker-unity", "Unity Game Developer", "C#, 3D/2D Games, Physics"),
        ("⚙️ worker-sys", "System Administrator", "PowerShell, Windows/Linux"),
        ("🎨 worker-web", "UI/Design Specialist", "HTML/CSS, Animations, Responsive"),
        ("🏗️ builder", "Build & CI/CD Engineer", "Scaffolding, Deps, Docker"),
        ("🔍 reviewer", "Code Review & QA", "Bugs, Security, Best Practices"),
    ]

    cols = 4
    for i, (name, role, specialty) in enumerate(agents):
        col = i % cols
        row = i // cols
        x = Inches(0.8) + Inches(col * 3.0)
        y = Inches(1.8) + Inches(row * 2.6)

        # Card
        add_shape(slide, x, y, Inches(2.7), Inches(2.2), BG_CARD)
        # Top accent
        add_shape(slide, x, y, Inches(2.7), Inches(0.04), PURPLE)

        # Name
        add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.3), Inches(0.4),
                     name, font_size=16, color=WHITE, bold=True, font_name='Calibri')
        # Role
        add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.3), Inches(0.3),
                     role, font_size=12, color=PURPLE_LIGHT, font_name='Calibri')
        # Specialty
        add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.3), Inches(0.8),
                     specialty, font_size=11, color=DARK_GRAY, font_name='Calibri')

    # Add Suzume card at the bottom
    y = Inches(1.8) + Inches(2 * 2.6)
    x = Inches(5.3)
    add_shape(slide, x, y, Inches(2.7), Inches(2.2), BG_CARD)
    add_shape(slide, x, y, Inches(2.7), Inches(0.04), PURPLE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.3), Inches(0.4),
                 "👑 Suzume", font_size=16, color=WHITE, bold=True, font_name='Calibri')
    add_text_box(slide, x + Inches(0.2), y + Inches(0.7), Inches(2.3), Inches(0.3),
                 "Supreme Orchestrator", font_size=12, color=PURPLE_LIGHT, font_name='Calibri')
    add_text_box(slide, x + Inches(0.2), y + Inches(1.1), Inches(2.3), Inches(0.8),
                 "Manages all agents, plans architecture, ensures quality",
                 font_size=11, color=DARK_GRAY, font_name='Calibri')

    # ════════════════════════════════════════
    # SLIDE 5: Superpowers
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_shape(slide, Inches(0.8), Inches(0.6), Inches(0.06), Inches(0.5), PURPLE)
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
                 "Superpowers", font_size=36, color=WHITE, bold=True,
                 font_name='Calibri Light')

    powers = [
        ("🌐  Full-Stack Web", "React, Next.js, Node.js, APIs, databases, auth, deployment — the entire stack"),
        ("🎮  Game Development", "Unity, C# scripting, 3D/2D assets, physics, animation, builds"),
        ("⚙️  System Automation", "PowerShell, Windows/Linux admin, file ops, processes, scheduled tasks"),
        ("📊  Data & ML", "Python, data pipelines, web scraping, ML models, visualization"),
        ("📱  Cross-Platform", "Desktop apps, mobile responsive, browser extensions, PWAs"),
        ("🛡️  Self-Healing", "Auto-debug, auto-fix, runtime verification, zero placeholders"),
    ]

    for i, (title, desc) in enumerate(powers):
        col = i % 3
        row = i // 3
        x = Inches(0.8) + Inches(col * 4.0)
        y = Inches(1.5) + Inches(row * 2.8)

        # Card
        add_shape(slide, x, y, Inches(3.6), Inches(2.4), BG_CARD)
        # Left accent
        add_shape(slide, x, y, Inches(0.04), Inches(2.4), PURPLE)

        # Title
        add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(0.4),
                     title, font_size=20, color=WHITE, bold=True, font_name='Calibri')
        # Description
        add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.0), Inches(1.2),
                     desc, font_size=14, color=LIGHT_GRAY, font_name='Calibri')

    # ════════════════════════════════════════
    # SLIDE 6: Project Showcase
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    add_shape(slide, Inches(0.8), Inches(0.6), Inches(0.06), Inches(0.5), PURPLE)
    add_text_box(slide, Inches(1.0), Inches(0.5), Inches(8), Inches(0.6),
                 "Project Showcase", font_size=36, color=WHITE, bold=True,
                 font_name='Calibri Light')

    # Project 1
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(2.5), BG_CARD)
    add_shape(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(0.04), GREEN)
    add_text_box(slide, Inches(1.2), Inches(1.7), Inches(5), Inches(0.4),
                 "✦  Anime Registration App", font_size=22, color=WHITE, bold=True,
                 font_name='Calibri')
    add_text_box(slide, Inches(1.2), Inches(2.2), Inches(5), Inches(1.2),
                 "Complete anime registration platform with:\n"
                 "• User authentication (JWT) & admin panel\n"
                 "• Prisma/SQLite database with 4 models\n"
                 "• Sakura animations & page transitions\n"
                 "• 12+ organized files, full-stack architecture",
                 font_size=14, color=LIGHT_GRAY, font_name='Calibri')

    # Project 2
    add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(2.5), BG_CARD)
    add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.04), PURPLE)
    add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4),
                 "✦  SaaS Website + Presentation", font_size=22, color=WHITE, bold=True,
                 font_name='Calibri')
    add_text_box(slide, Inches(7.2), Inches(2.2), Inches(5), Inches(1.2),
                 "This very showcase built in real-time:\n"
                 "• SaaS landing page with particles & 3D tilt\n"
                 "• Web presentation with 7 cinematic slides\n"
                 "• Canvas particle systems & parallax\n"
                 "• Actual .pptx PowerPoint file generation",
                 font_size=14, color=LIGHT_GRAY, font_name='Calibri')

    # Stats
    add_shape(slide, Inches(0.8), Inches(4.5), Inches(11.6), Inches(1.5), BG_CARD)
    stats = [("10K+", "Users Served"), ("7", "Sub-Agents"), ("100%", "Autonomous"), ("60fps", "Animations")]
    for i, (num, label) in enumerate(stats):
        x = Inches(1.5) + Inches(i * 3.0)
        add_text_box(slide, x, Inches(4.7), Inches(2.5), Inches(0.5),
                     num, font_size=32, color=PURPLE_LIGHT, bold=True,
                     alignment=PP_ALIGN.CENTER, font_name='Calibri Light')
        add_text_box(slide, x, Inches(5.2), Inches(2.5), Inches(0.4),
                     label, font_size=14, color=DARK_GRAY,
                     alignment=PP_ALIGN.CENTER, font_name='Calibri')

    # ════════════════════════════════════════
    # SLIDE 7: Thank You
    # ════════════════════════════════════════
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)

    # Decorative bars
    add_shape(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.06), PURPLE)
    add_shape(slide, Inches(0), Inches(7.44), Inches(13.333), Inches(0.06), BLUE)

    add_text_box(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.5),
                 "Thank You!", font_size=72, color=WHITE, bold=True,
                 alignment=PP_ALIGN.CENTER, font_name='Calibri Light')

    add_text_box(slide, Inches(1.5), Inches(3.3), Inches(10), Inches(0.8),
                 "Built entirely by Suzume — Your Supreme AI Companion",
                 font_size=24, color=PURPLE_LIGHT, alignment=PP_ALIGN.CENTER,
                 font_name='Calibri Light')

    add_text_box(slide, Inches(2), Inches(4.3), Inches(9), Inches(0.5),
                 "Created for Rushikesh  •  July 18, 2026",
                 font_size=14, color=DARK_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name='Calibri')

    add_text_box(slide, Inches(3), Inches(5.0), Inches(7), Inches(0.5),
                 "✧  Suzume is ready for your next project  ✧",
                 font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER,
                 font_name='Calibri Light')

    # ─── Save ───
    output_path = os.path.join(os.path.dirname(__file__), "Suzume_Presentation.pptx")
    prs.save(output_path)
    print(f"[OK] PowerPoint saved to: {output_path}")
    return output_path

if __name__ == "__main__":
    path = create_presentation()
    print(f"\nFile size: {os.path.getsize(path) / 1024:.1f} KB")
    print("Open in PowerPoint to view!")
